"""FoodFact v4 — présentation visuelle enrichie (~15 min + démo live)."""

import os
import sys
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from presentation_mock_ui import (
    mock_architecture_diagram,
    mock_browser_frame,
    mock_chart_bars,
    mock_data_table,
    mock_filters_panel,
    mock_json_snippet,
    mock_product_card,
    mock_search_bar,
    mock_terminal_block,
    NUTRI,
    EMERALD,
    CYAN,
    VIOLET,
    TEXT_DARK,
    TEXT_MUTED,
    WHITE,
    CARD_BG,
    CARD_BORDER,
    BG_PANEL,
)

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT = os.path.join(BASE, "PRESENTATION_FoodFact_v4.pptx")

# ── Theme ────────────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0B, 0x11, 0x20)
BG_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = EMERALD
ACCENT2 = CYAN
ACCENT3 = VIOLET
FOOTER_DARK = RGBColor(0x94, 0xA3, 0xB8)
FONT = "Segoe UI"

W, H = Inches(10), Inches(7.5)
TOTAL_SLIDES = 18


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _style_paragraph(p, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = FONT
    p.font.color.rgb = color


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def circle(slide, l, t, size, fill, transparency=0.0):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.fill.transparency = transparency
    c.line.fill.background()
    return c


def send_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(2, el)


def textbox(slide, l, t, w, h, text, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _style_paragraph(p, size, bold, color, align)
    return box


def multiline(slide, l, t, w, h, lines, size=16, color=TEXT_DARK, spacing=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        _style_paragraph(p, size, False, color)
        p.space_after = Pt(spacing)
    return box


def bullets(slide, l, t, w, h, items, size=14, color=TEXT_DARK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        _style_paragraph(p, size, False, color)
        p.space_after = Pt(6)
    return box


def slide_bg(slide, dark=False):
    bg = rect(slide, Inches(0), Inches(0), W, H, BG_DARK if dark else BG_LIGHT)
    send_to_back(slide, bg)


def accent_bar(slide):
    rect(slide, Inches(0), Inches(0), Inches(0.1), H, ACCENT)


def footer(slide, num, dark=False, label="FoodFact · NutriRecherche"):
    fc = FOOTER_DARK if dark else TEXT_MUTED
    textbox(slide, Inches(0.55), Inches(7.05), Inches(5), Inches(0.3), label, 9, False, fc)
    textbox(slide, Inches(8.8), Inches(7.05), Inches(0.8), Inches(0.3), f"{num}/{TOTAL_SLIDES}", 9, False, fc, PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None, num=0, dark=False):
    accent_bar(slide)
    tc = WHITE if dark else TEXT_DARK
    sc = ACCENT if dark else TEXT_MUTED
    textbox(slide, Inches(0.55), Inches(0.4), Inches(8.5), Inches(0.75), title, 32, True, tc)
    if subtitle:
        textbox(slide, Inches(0.55), Inches(1.0), Inches(8.5), Inches(0.4), subtitle, 14, False, sc)
    if num:
        footer(slide, num, dark)


def card(slide, l, t, w, h, title, body=None, accent=ACCENT):
    rect(slide, l, t, w, h, CARD_BG, CARD_BORDER)
    rect(slide, l, t, w, Inches(0.055), accent)
    textbox(slide, l + Inches(0.18), t + Inches(0.18), w - Inches(0.36), Inches(0.42), title, 14, True, TEXT_DARK)
    if body:
        if isinstance(body, str):
            body = [body]
        bullets(slide, l + Inches(0.18), t + Inches(0.58), w - Inches(0.36), h - Inches(0.65), body, 12, TEXT_MUTED)


def pill_row(slide, top, labels, colors=None, left=Inches(0.55)):
    colors = colors or [ACCENT] * len(labels)
    x = left
    for label, col in zip(labels, colors):
        w = Inches(max(1.0, len(label) * 0.105 + 0.45))
        sh = rect(slide, x, top, w, Inches(0.38), col)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = WHITE
        x += w + Inches(0.12)


def stat_block(slide, l, t, w, value, label, color=ACCENT):
    rect(slide, l, t, w, Inches(1.25), CARD_BG, CARD_BORDER)
    rect(slide, l, t, w, Inches(0.05), color)
    textbox(slide, l, t + Inches(0.18), w, Inches(0.6), value, 26, True, color, PP_ALIGN.CENTER)
    textbox(slide, l, t + Inches(0.78), w, Inches(0.35), label, 11, False, TEXT_MUTED, PP_ALIGN.CENTER)


def section_divider(slide, title, subtitle, num):
    slide_bg(slide, dark=True)
    accent_bar(slide)
    circle(slide, Inches(7.5), Inches(-1.2), Inches(4), ACCENT, 0.88)
    circle(slide, Inches(-1.5), Inches(5.5), Inches(3.5), ACCENT2, 0.9)
    textbox(slide, Inches(0.7), Inches(2.8), Inches(8.6), Inches(1), title, 44, True, WHITE, PP_ALIGN.CENTER)
    textbox(slide, Inches(1.2), Inches(3.85), Inches(7.6), Inches(0.6), subtitle, 18, False, ACCENT, PP_ALIGN.CENTER)
    footer(slide, num, dark=True)


def icon_dot(slide, l, t, color, size=Inches(0.14)):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c


def numbered_step(slide, l, t, num, title, desc, color=ACCENT):
    rect(slide, l, t, Inches(8.9), Inches(0.85), CARD_BG, CARD_BORDER)
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.15), t + Inches(0.18), Inches(0.48), Inches(0.48))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE
    textbox(slide, l + Inches(0.8), t + Inches(0.12), Inches(3), Inches(0.35), title, 13, True, TEXT_DARK)
    textbox(slide, l + Inches(0.8), t + Inches(0.45), Inches(8), Inches(0.35), desc, 11, False, TEXT_MUTED)


# ─── SLIDES ─────────────────────────────────────────────────────────────────

def s01_cover(prs):
    slide = blank(prs)
    slide_bg(slide, dark=True)
    accent_bar(slide)
    circle(slide, Inches(6.8), Inches(-1), Inches(4.2), ACCENT, 0.82)
    circle(slide, Inches(8.2), Inches(4.5), Inches(2.8), ACCENT2, 0.88)
    circle(slide, Inches(-0.8), Inches(5.2), Inches(3), ACCENT3, 0.9)

    pill_row(slide, Inches(0.7), ["Open Data", "Nutrition", "Next.js + Scala"], [ACCENT, ACCENT2, ACCENT3], Inches(0.7))

    textbox(slide, Inches(0.7), Inches(2.0), Inches(8.5), Inches(1.1), "FoodFact", 58, True, WHITE)
    textbox(slide, Inches(0.7), Inches(3.15), Inches(8.5), Inches(0.55),
            "Explorer · Comparer · Manger mieux", 22, False, ACCENT)
    textbox(slide, Inches(0.7), Inches(3.85), Inches(8.5), Inches(0.5),
            "Application web alimentée par l'Open Data d'OpenFoodFacts", 15, False, FOOTER_DARK)

    multiline(slide, Inches(0.7), Inches(5.3), Inches(8.5), Inches(0.8),
              [f"NutriRecherche  ·  {date.today().strftime('%d/%m/%Y')}  ·  MIT License"],
              12, FOOTER_DARK, 4)
    footer(slide, 1, dark=True)


def s02_agenda(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Plan de la présentation", "15 minutes · Open Data, technique & démo live (~5 min)", 2)

    items = [
        ("01", "3 min", "Contexte & problématique", "Pourquoi FoodFact ?"),
        ("02", "3 min", "Open Data & OpenFoodFacts", "Source, licence, réutilisation"),
        ("03", "4 min", "Architecture & stack", "Frontend, backend, flux de données"),
        ("04", "3 min", "Fonctionnalités & mockups UI", "Ce que voit l'utilisateur"),
        ("05", "5 min", "Démo live + conclusion", "Lancement app · scénario yaourt"),
    ]
    y = Inches(1.55)
    for num, timing, title, sub in items:
        rect(slide, Inches(0.55), y, Inches(8.9), Inches(0.95), CARD_BG, CARD_BORDER)
        rect(slide, Inches(0.55), y, Inches(0.55), Inches(0.95), ACCENT if num in ("01", "03", "05") else ACCENT2)
        textbox(slide, Inches(0.65), y + Inches(0.12), Inches(0.45), Inches(0.5), num, 16, True, WHITE, PP_ALIGN.CENTER)
        textbox(slide, Inches(1.25), y + Inches(0.1), Inches(1.0), Inches(0.35), timing, 11, True, ACCENT2)
        textbox(slide, Inches(2.3), y + Inches(0.08), Inches(4), Inches(0.4), title, 14, True, TEXT_DARK)
        textbox(slide, Inches(2.3), y + Inches(0.45), Inches(6.5), Inches(0.35), sub, 11, False, TEXT_MUTED)
        y += Inches(1.05)


def s03_constat(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Le constat", "Avant / Après FoodFact", 3)

    # Avant
    rect(slide, Inches(0.55), Inches(1.55), Inches(4.2), Inches(4.8), RGBColor(0xFE, 0xF2, 0xF2), CARD_BORDER)
    rect(slide, Inches(0.55), Inches(1.55), Inches(4.2), Inches(0.06), RGBColor(0xEF, 0x44, 0x44))
    textbox(slide, Inches(0.75), Inches(1.75), Inches(3.8), Inches(0.35), "Sans FoodFact", 16, True, RGBColor(0xDC, 0x26, 0x26))
    bullets(slide, Inches(0.75), Inches(2.2), Inches(3.8), Inches(3.5), [
        "Étiquettes illisibles en magasin",
        "Scores multiples difficiles à comparer",
        "Pas de vue d'ensemble sur une catégorie",
        "Recherche manuelle produit par produit",
    ], 12, TEXT_MUTED)

    # Après
    rect(slide, Inches(5.25), Inches(1.55), Inches(4.2), Inches(4.8), RGBColor(0xEC, 0xFD, 0xF5), CARD_BORDER)
    rect(slide, Inches(5.25), Inches(1.55), Inches(4.2), Inches(0.06), ACCENT)
    textbox(slide, Inches(5.45), Inches(1.75), Inches(3.8), Inches(0.35), "Avec FoodFact", 16, True, ACCENT)
    bullets(slide, Inches(5.45), Inches(2.2), Inches(3.8), Inches(3.5), [
        "Recherche instantanée sur 3M+ produits",
        "Filtres Nutri-Score, NOVA, nutriments",
        "Comparaison grille / tableau / graphiques",
        "Alternatives plus saines en un clic",
    ], 12, TEXT_MUTED)

    textbox(slide, Inches(0.55), Inches(6.55), Inches(8.9), Inches(0.35),
            "Objectif : rendre l'information nutritionnelle accessible, comparable et actionnable.",
            12, True, TEXT_MUTED, PP_ALIGN.CENTER)


def s04_solution(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Notre réponse", "Problème → Solution → Résultat", 4)

    cards = [
        ("Problème", ["Info nutritionnelle fragmentée et complexe", "Comparaison difficile entre produits"], RGBColor(0xEF, 0x44, 0x44)),
        ("Solution", ["Réutiliser l'Open Data OpenFoodFacts", "Backend Scala + Frontend Next.js"], ACCENT),
        ("Résultat", ["UX claire · filtres · scores visuels", "URL partageable · mode clair/sombre"], ACCENT3),
    ]
    x = Inches(0.55)
    for title, body, col in cards:
        card(slide, x, Inches(1.65), Inches(2.85), Inches(2.2), title, body, col)
        x += Inches(3.05)

    rect(slide, Inches(0.55), Inches(4.2), Inches(8.9), Inches(2.2), CARD_BG, CARD_BORDER)
    rect(slide, Inches(0.55), Inches(4.2), Inches(8.9), Inches(0.05), ACCENT2)
    textbox(slide, Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.35), "Proposition de valeur", 15, True, TEXT_DARK)
    bullets(slide, Inches(0.8), Inches(4.85), Inches(8.2), Inches(1.4), [
        "Consommateur d'Open Data — pas de base propriétaire à maintenir",
        "Couche API intelligente : cache, rate limit, filtres, normalisation pays",
        "Interface moderne pour explorer, comparer et choisir mieux",
    ], 12, TEXT_MUTED)


def s05_open_data(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Qu'est-ce que l'Open Data ?", "Données publiques, ouvertes et réutilisables", 5)

    defs = [
        ("Gratuit", "Accès sans frais", ACCENT),
        ("Ouvert", "JSON / API REST", ACCENT2),
        ("Réutilisable", "Licence ODbL", ACCENT3),
        ("Impact", "Apps citoyennes", ACCENT),
    ]
    x = Inches(0.55)
    for title, body, col in defs:
        card(slide, x, Inches(1.55), Inches(2.05), Inches(1.35), title, [body], col)
        x += Inches(2.2)

    # Producer vs consumer diagram
    rect(slide, Inches(0.55), Inches(3.15), Inches(4.2), Inches(3.5), CARD_BG, CARD_BORDER)
    textbox(slide, Inches(0.75), Inches(3.3), Inches(3.8), Inches(0.3), "Producteur", 13, True, ACCENT3)
    bullets(slide, Inches(0.75), Inches(3.65), Inches(3.8), Inches(2.5), [
        "OpenFoodFacts collecte les données",
        "Communauté mondiale (crowdsourcing)",
        "Licence ODbL — partage autorisé",
    ], 11, TEXT_MUTED)

    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.85), Inches(4.5), Inches(0.5), Inches(0.35))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT
    arr.line.fill.background()

    rect(slide, Inches(5.5), Inches(3.15), Inches(3.95), Inches(3.5), CARD_BG, CARD_BORDER)
    rect(slide, Inches(5.5), Inches(3.15), Inches(3.95), Inches(0.06), ACCENT)
    textbox(slide, Inches(5.7), Inches(3.3), Inches(3.5), Inches(0.3), "Consommateur — FoodFact", 13, True, ACCENT)
    bullets(slide, Inches(5.7), Inches(3.65), Inches(3.5), Inches(2.5), [
        "Valorise via UX et API",
        "Ne produit pas les données",
        "Transparence · réutilisation · open source",
    ], 11, TEXT_MUTED)


def s06_off(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "OpenFoodFacts", "world.openfoodfacts.org", 6)

    stat_block(slide, Inches(0.55), Inches(1.45), Inches(2.0), "3M+", "produits", ACCENT)
    stat_block(slide, Inches(2.7), Inches(1.45), Inches(2.0), "180+", "pays", ACCENT2)
    stat_block(slide, Inches(4.85), Inches(1.45), Inches(2.0), "ODbL", "licence", ACCENT3)
    stat_block(slide, Inches(7.0), Inches(1.45), Inches(2.45), "API", "REST publique", ACCENT)

    card(slide, Inches(0.55), Inches(2.95), Inches(4.2), Inches(2.8),
         "Comment ça marche ?",
         ["Scans code-barres via app mobile", "Contributions bénévoles mondiales",
          "Scores, ingrédients, images enrichis", "Mise à jour continue"])

    card(slide, Inches(5.0), Inches(2.95), Inches(4.45), Inches(2.8),
         "Données exploitées",
         ["Nutri-Score · Éco-Score · NOVA", "Nutriments / 100g",
          "Ingrédients · allergènes · marques", "Images + lien fiche OFF"],
         ACCENT2)

    x = Inches(0.55)
    for grade, col in NUTRI.items():
        sh = rect(slide, x, Inches(6.0), Inches(0.5), Inches(0.3), col)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = grade
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = WHITE
        x += Inches(0.58)
    textbox(slide, Inches(3.5), Inches(5.95), Inches(5.5), Inches(0.35),
            "Échelle Nutri-Score (A = meilleur · E = moins bon)", 10, False, TEXT_MUTED)


def s07_open_data_oui(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Open Data confirmé", "OpenFoodFacts = source ouverte, gratuite, réutilisable", 7)

    mock_json_snippet(slide, Inches(0.55), Inches(1.55), Inches(4.3), Inches(2.8), [
        '{ "code": "3560070460096",',
        '  "product_name": "Yaourt nature",',
        '  "brands": "Danone",',
        '  "nutriscore_grade": "a",',
        '  "nova_group": 1,',
        '  "nutriments": { "energy-kcal_100g": 61 }',
        '}',
    ])

    rect(slide, Inches(5.1), Inches(1.55), Inches(4.35), Inches(2.8), CARD_BG, CARD_BORDER)
    rect(slide, Inches(5.1), Inches(1.55), Inches(4.35), Inches(0.05), ACCENT)
    textbox(slide, Inches(5.3), Inches(1.75), Inches(4), Inches(0.35), "Pourquoi c'est de l'Open Data", 14, True, TEXT_DARK)
    bullets(slide, Inches(5.3), Inches(2.15), Inches(3.9), Inches(2.0), [
        "API REST publique sans clé API",
        "Format JSON ouvert et documenté",
        "Licence ODbL — réutilisation autorisée",
        "FoodFact consomme, ne verrouille pas",
    ], 11, TEXT_MUTED)

    textbox(slide, Inches(0.55), Inches(4.6), Inches(8.9), Inches(0.4),
            "Chaîne : OpenFoodFacts  →  Backend Scala  →  Frontend Next.js  →  Utilisateur",
            13, True, ACCENT, PP_ALIGN.CENTER)


def s08_architecture(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Architecture", "Client-serveur · 3 couches · séparation des responsabilités", 8)

    mock_architecture_diagram(slide, Inches(0.55), Inches(1.65), Inches(8.9), Inches(2.2))

    layers = [
        ("Frontend — Port 3000", "UI · Filtres · Graphiques · URL partageable · thème clair/sombre", ACCENT),
        ("Backend — Port 8080", "Proxy OFF · Cache 300s · Rate limit 30/min · Filtres & tri", ACCENT2),
        ("OpenFoodFacts — Externe", "Source de vérité · 3M+ fiches JSON · Pas de lock-in propriétaire", ACCENT3),
    ]
    y = Inches(4.15)
    for title, desc, col in layers:
        icon_dot(slide, Inches(0.65), y + Inches(0.08), col)
        textbox(slide, Inches(0.9), y, Inches(8.3), Inches(0.3), title, 12, True, TEXT_DARK)
        textbox(slide, Inches(0.9), y + Inches(0.28), Inches(8.3), Inches(0.3), desc, 10, False, TEXT_MUTED)
        y += Inches(0.75)


def s09_flux(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Flux de données", "Requête GET /api/search — de l'utilisateur à OpenFoodFacts", 9)

    steps = [
        (1, "Utilisateur", "Saisit « yaourt » + filtre Nutri-Score A/B", ACCENT),
        (2, "Frontend", "Debounce 300 ms · sync URL · appel API", ACCENT2),
        (3, "Backend", "Vérifie cache · rate limit · construit requête OFF", ACCENT3),
        (4, "OpenFoodFacts", "Retourne JSON produits bruts", ACCENT),
        (5, "Affichage", "Filtres appliqués · tri · pagination · ProductCard", ACCENT2),
    ]
    y = Inches(1.55)
    for num, title, desc, col in steps:
        numbered_step(slide, Inches(0.55), y, num, title, desc, col)
        y += Inches(0.95)


def s10_stack(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Stack technique", "Choix motivés par typage fort, effets fonctionnels et Open Data", 10)

    sections = [
        ("Frontend", ACCENT, ["Next.js 14+", "React", "TypeScript", "Tailwind", "shadcn/ui", "Recharts"]),
        ("Backend", ACCENT2, ["Scala 3.3", "http4s", "Circe", "Cats Effect", "Ember", "MUnit"]),
        ("Ops & Open Data", ACCENT3, ["OpenFoodFacts", "Docker", "Git", "CORS", "Cache", "Rate limit"]),
    ]
    y = Inches(1.5)
    for label, col, pills in sections:
        rect(slide, Inches(0.55), y, Inches(8.9), Inches(1.35), CARD_BG, CARD_BORDER)
        rect(slide, Inches(0.55), y, Inches(8.9), Inches(0.05), col)
        textbox(slide, Inches(0.75), y + Inches(0.12), Inches(3), Inches(0.3), label, 14, True, col)
        pill_row(slide, y + Inches(0.52), pills, [col] * len(pills))
        y += Inches(1.48)

    rect(slide, Inches(0.55), Inches(5.95), Inches(8.9), Inches(0.85), CARD_BG, CARD_BORDER)
    textbox(slide, Inches(0.8), Inches(6.1), Inches(8.4), Inches(0.55),
            "Pourquoi Scala + Next.js ? Séparation claire · typage TS + Scala · effets Cats · écosystème React riche",
            11, False, TEXT_MUTED, PP_ALIGN.CENTER)


def s11_features(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Fonctionnalités", "Parcours utilisateur complet", 11)

    feats = [
        ("Recherche", "Debounce 300 ms · URL partageable", ACCENT),
        ("Filtres", "Marque · Pays · Nutri · NOVA · nutriments", ACCENT2),
        ("Tri & pages", "Énergie, sucres, graisses · Charger plus", ACCENT3),
        ("3 vues", "Grille · Tableau · Statistiques", ACCENT),
        ("Fiche produit", "Ingrédients · Allergènes · Lien OFF", ACCENT2),
        ("Alternatives", "Similaires triés par Nutri-Score", ACCENT3),
    ]
    x, y = Inches(0.55), Inches(1.55)
    for i, (t, b, col) in enumerate(feats):
        col_i, row = i % 3, i // 3
        card(slide, x + col_i * Inches(3.05), y + row * Inches(1.95), Inches(2.85), Inches(1.65), t, [b], col)

    textbox(slide, Inches(0.55), Inches(5.75), Inches(8.9), Inches(0.35),
            "yaourt → Nutri A/B → tableau → fiche → alternatives", 12, True, ACCENT, PP_ALIGN.CENTER)


def s12_mock_search(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Interface — Recherche", "Mockup UI · page d'accueil", 12)

    cx, cy, cw, ch = mock_browser_frame(
        slide, Inches(0.55), Inches(1.5), Inches(8.9), Inches(5.3),
        "localhost:3000/?q=yaourt&nutriscore=a,b"
    )

    mock_filters_panel(slide, cx + Inches(0.1), cy + Inches(0.1), Inches(1.85), ch - Inches(0.2))

    content_x = cx + Inches(2.05)
    content_w = cw - Inches(2.15)
    mock_search_bar(slide, content_x, cy + Inches(0.1), content_w, "yaourt", "grid")

    card_w = (content_w - Inches(0.15)) / 2
    mock_product_card(slide, content_x, cy + Inches(0.65), card_w, Inches(2.5),
                      "Yaourt nature", "A", "NOVA 1", "Danone", 0)
    mock_product_card(slide, content_x + card_w + Inches(0.15), cy + Inches(0.65), card_w, Inches(2.5),
                      "Yaourt aux fruits", "B", "NOVA 2", "Yoplait", 1)


def s13_mock_analyse(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Interface — Analyse", "Mockup UI · vues Tableau & Statistiques", 13)

    cx, cy, cw, ch = mock_browser_frame(
        slide, Inches(0.55), Inches(1.5), Inches(8.9), Inches(5.3),
        "localhost:3000/?q=yaourt&view=stats"
    )

    mock_data_table(slide, cx + Inches(0.15), cy + Inches(0.15), cw * 0.55, ch - Inches(0.3))
    mock_chart_bars(slide, cx + cw * 0.58, cy + Inches(0.15), cw * 0.38, ch - Inches(0.3))


def s14_api(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Notre API REST", "Couche intelligente entre l'UI et OpenFoodFacts", 14)

    endpoints = [
        ("GET /health", '{"status":"ok"}', ACCENT),
        ("GET /api/search", "12 paramètres · filtres · tri · pagination", ACCENT2),
        ("GET /api/product/{code}", "Détail + alternatives par catégorie", ACCENT3),
    ]
    y = Inches(1.5)
    for path, desc, col in endpoints:
        rect(slide, Inches(0.55), y, Inches(4.2), Inches(0.75), CARD_BG, CARD_BORDER)
        rect(slide, Inches(0.55), y, Inches(0.06), Inches(0.75), col)
        textbox(slide, Inches(0.75), y + Inches(0.1), Inches(3.8), Inches(0.3), path, 12, True, col)
        textbox(slide, Inches(0.75), y + Inches(0.38), Inches(3.8), Inches(0.3), desc, 10, False, TEXT_MUTED)
        y += Inches(0.85)

    mock_json_snippet(slide, Inches(5.0), Inches(1.5), Inches(4.45), Inches(2.55), [
        '{ "count": 12, "totalFromOff": 1234,',
        '  "page": 1, "pageSize": 50,',
        '  "products": [',
        '    { "code": "...", "product_name": "...",',
        '      "nutriscore_grade": "a" }',
        '  ]',
        '}',
    ])

    params = "q · brand · country · nutriscore · nova · minEnergy · maxEnergy · minSugar · maxSugar · minFat · maxFat · sortBy · order · page · pageSize"
    textbox(slide, Inches(0.55), Inches(4.15), Inches(8.9), Inches(0.55), f"Paramètres /api/search : {params}", 9, False, TEXT_MUTED)

    cards = [
        ("Cache", "TTL 300 s", ACCENT),
        ("Rate limit", "30 req/min/IP", ACCENT2),
        ("Pays", "CountryUtils", ACCENT3),
        ("Erreurs", "502 · 504 · 400", ACCENT),
    ]
    x = Inches(0.55)
    for title, body, col in cards:
        card(slide, x, Inches(4.85), Inches(2.05), Inches(1.15), title, [body], col)
        x += Inches(2.2)


def s15_defis(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Défis & solutions", "Problèmes rencontrés et réponses techniques", 15)

    defis = [
        ("Latence OpenFoodFacts", "Timeout configurable (OFF_TIMEOUT_MS=10s) · erreur 504", ACCENT),
        ("Surcharge API externe", "Cache mémoire 300s · rate limit 30 req/min/IP", ACCENT2),
        ("Pays hétérogènes", "CountryUtils — normalisation noms et codes pays", ACCENT3),
        ("Erreurs upstream", "502 Bad Gateway · 504 Gateway Timeout · messages JSON clairs", ACCENT),
    ]
    x, y = Inches(0.55), Inches(1.55)
    for i, (title, desc, col) in enumerate(defis):
        col_i, row = i % 2, i // 2
        rect(slide, x + col_i * Inches(4.55), y + row * Inches(2.05), Inches(4.2), Inches(1.75), CARD_BG, CARD_BORDER)
        rect(slide, x + col_i * Inches(4.55), y + row * Inches(2.05), Inches(4.2), Inches(0.05), col)
        textbox(slide, x + col_i * Inches(4.55) + Inches(0.2), y + row * Inches(2.05) + Inches(0.15),
                Inches(3.8), Inches(0.35), title, 13, True, TEXT_DARK)
        textbox(slide, x + col_i * Inches(4.55) + Inches(0.2), y + row * Inches(2.05) + Inches(0.52),
                Inches(3.8), Inches(0.9), desc, 11, False, TEXT_MUTED)


def s16_demo_live(prs):
    slide = blank(prs)
    slide_bg(slide, dark=True)
    accent_bar(slide)
    header(slide, "Démo live", "Lancement · scénario oral · ~5 minutes", 16, dark=True)

    mock_terminal_block(slide, Inches(0.55), Inches(1.55), Inches(4.2), Inches(2.5), [
        "$ cd backend_scala",
        "$ .\\sbt.bat run",
        "# → http://localhost:8080/health",
        "",
        "$ cd frontend_react",
        "$ npm run dev",
        "# → http://localhost:3000",
    ])

    rect(slide, Inches(5.0), Inches(1.55), Inches(4.45), Inches(4.8), RGBColor(0x15, 0x23, 0x42))
    textbox(slide, Inches(5.2), Inches(1.7), Inches(4.1), Inches(0.35), "Scénario yaourt (5 étapes)", 14, True, ACCENT)
    scenario = [
        "1. Rechercher « yaourt »",
        "2. Filtrer Nutri-Score A et B",
        "3. Passer en vue Tableau",
        "4. Ouvrir une fiche produit",
        "5. Montrer alternatives + lien OFF",
    ]
    y = Inches(2.15)
    for step in scenario:
        textbox(slide, Inches(5.2), y, Inches(4.1), Inches(0.4), step, 12, False, FOOTER_DARK)
        y += Inches(0.48)

    textbox(slide, Inches(0.55), Inches(6.55), Inches(8.9), Inches(0.35),
            "Backend :8080  ·  Frontend :3000  ·  URL partageable  ·  thème clair/sombre",
            11, False, FOOTER_DARK, PP_ALIGN.CENTER)
    footer(slide, 16, dark=True)


def s17_perspectives(prs):
    slide = blank(prs)
    slide_bg(slide)
    header(slide, "Perspectives", "Évolutions possibles", 17)

    items = [
        ("Favoris", "Sauvegarder des produits en local ou compte utilisateur", ACCENT),
        ("Comparaison", "Vue side-by-side de 2 à 4 produits", ACCENT2),
        ("Export CSV", "Télécharger les résultats filtrés", ACCENT3),
        ("PWA", "Application installable sur mobile", ACCENT),
    ]
    x, y = Inches(0.55), Inches(1.55)
    for i, (title, desc, col) in enumerate(items):
        col_i, row = i % 2, i // 2
        card(slide, x + col_i * Inches(4.55), y + row * Inches(2.05), Inches(4.2), Inches(1.65), title, [desc], col)


def s18_conclusion(prs):
    slide = blank(prs)
    slide_bg(slide, dark=True)
    accent_bar(slide)
    circle(slide, Inches(7.8), Inches(-0.5), Inches(3.5), ACCENT, 0.85)
    circle(slide, Inches(-1.2), Inches(5.8), Inches(3.8), ACCENT2, 0.88)

    textbox(slide, Inches(0.7), Inches(1.4), Inches(8.6), Inches(0.8), "En résumé", 38, True, WHITE, PP_ALIGN.CENTER)

    points = [
        ("Open Data", "OpenFoodFacts — 3M+ produits, licence ODbL, API publique"),
        ("Application", "FoodFact valorise les données via UX, filtres et comparaisons"),
        ("Architecture", "Next.js + Scala/http4s — cache, rate limit, tests MUnit"),
        ("Impact", "Manger mieux grâce à des données publiques réutilisables"),
    ]
    y = Inches(2.35)
    for title, desc in points:
        rect(slide, Inches(1.2), y, Inches(7.6), Inches(0.8), RGBColor(0x15, 0x23, 0x42))
        textbox(slide, Inches(1.45), y + Inches(0.1), Inches(2.2), Inches(0.35), title, 13, True, ACCENT)
        textbox(slide, Inches(3.6), y + Inches(0.12), Inches(5), Inches(0.55), desc, 12, False, FOOTER_DARK)
        y += Inches(0.95)

    textbox(slide, Inches(0.7), Inches(6.0), Inches(8.6), Inches(0.5),
            "Merci — Questions ?", 26, True, WHITE, PP_ALIGN.CENTER)
    textbox(slide, Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.35),
            "github.com/Yamnyr/NutriRecherche", 11, False, ACCENT, PP_ALIGN.CENTER)
    footer(slide, 18, dark=True)


def save_presentation(prs):
    try:
        prs.save(OUTPUT)
        return OUTPUT
    except PermissionError:
        import time
        alt = OUTPUT.replace(".pptx", f"_{int(time.time())}.pptx")
        prs.save(alt)
        return alt


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    builders = [
        s01_cover, s02_agenda, s03_constat, s04_solution, s05_open_data,
        s06_off, s07_open_data_oui, s08_architecture, s09_flux, s10_stack,
        s11_features, s12_mock_search, s13_mock_analyse, s14_api,
        s15_defis, s16_demo_live, s17_perspectives, s18_conclusion,
    ]
    for fn in builders:
        fn(prs)

    path = save_presentation(prs)
    print(f"OK: {path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

"""Convertit FoodFact Presentation.pdf en PPTX avec contenu centré."""

import io
import os
import time

import fitz
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PDF_PATH = os.path.join(BASE, "FoodFact Presentation.pdf")
OUTPUT_PATH = os.path.join(BASE, "FoodFact Presentation.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TARGET_W, TARGET_H = 2560, 1440
FOOTER_H = 100
CONTENT_THRESHOLD = 180


def center_slide_image(img_bytes: bytes) -> bytes:
    """Recadre chaque slide pour centrer le contenu principal (hors footer)."""
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    arr = np.array(img, dtype=np.float32)
    gray = arr.mean(axis=2)

    work = gray.copy()
    work[-FOOTER_H:, :] = 255
    mask = work < CONTENT_THRESHOLD
    rows = np.where(mask.any(axis=1))[0]
    cols = np.where(mask.any(axis=0))[0]
    if len(rows) == 0 or len(cols) == 0:
        return img_bytes

    r0, r1 = rows[0], rows[-1]
    c0, c1 = cols[0], cols[-1]
    cy = (r0 + r1) / 2
    cx = (c0 + c1) / 2

    target_cy = (TARGET_H - FOOTER_H) / 2
    target_cx = TARGET_W / 2
    shift_y = int(target_cy - cy)
    shift_x = int(target_cx - cx)

    bg = tuple(np.mean([arr[0, 0], arr[50, 50], arr[0, 50]], axis=0).astype(np.uint8))
    canvas = Image.new("RGB", (TARGET_W, TARGET_H), bg)
    canvas.paste(img, (shift_x, shift_y))

    out = io.BytesIO()
    canvas.save(out, format="JPEG", quality=95)
    return out.getvalue()


def convert(pdf_path: str, output_path: str, center: bool = True) -> str:
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"PDF introuvable : {pdf_path}")

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for page in doc:
        slide = prs.slides.add_slide(blank)
        images = page.get_images(full=True)
        if not images:
            continue
        raw = doc.extract_image(images[0][0])["image"]
        data = center_slide_image(raw) if center else raw
        slide.shapes.add_picture(
            io.BytesIO(data),
            Inches(0),
            Inches(0),
            width=SLIDE_W,
            height=SLIDE_H,
        )

    doc.close()

    try:
        prs.save(output_path)
        return output_path
    except PermissionError:
        alt = output_path.replace(".pptx", f"_{int(time.time())}.pptx")
        prs.save(alt)
        return alt


def main():
    path = convert(PDF_PATH, OUTPUT_PATH, center=True)
    n = len(Presentation(path).slides)
    print(f"OK: {path} ({n} slides, contenu centré)")


if __name__ == "__main__":
    main()

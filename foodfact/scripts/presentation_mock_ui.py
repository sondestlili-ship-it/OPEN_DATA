"""Mock UI builders for FoodFact presentation slides."""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (aligned with frontend) ───────────────────────────────────────────
EMERALD = RGBColor(0x10, 0xB9, 0x81)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BG_PANEL = RGBColor(0xF8, 0xFA, 0xFC)
BG_CODE = RGBColor(0x0F, 0x17, 0x2A)
CODE_TEXT = RGBColor(0x94, 0xA3, 0xB8)
CODE_ACCENT = RGBColor(0x34, 0xD3, 0x99)

NUTRI = {
    "A": RGBColor(0x22, 0xC5, 0x5E),
    "B": RGBColor(0x84, 0xCC, 0x16),
    "C": RGBColor(0xEA, 0xB3, 0x08),
    "D": RGBColor(0xF9, 0x73, 0x16),
    "E": RGBColor(0xEF, 0x44, 0x44),
}

CARD_GRADIENTS = [
    (RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0xEC, 0xFD, 0xF5)),
    (RGBColor(0xCC, 0xFB, 0xF1), RGBColor(0xF0, 0xFD, 0xFA)),
]


def _rect(slide, l, t, w, h, fill, line=None, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(radius_shape, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _text(slide, l, t, w, h, text, size=11, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font="Segoe UI"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    return box


def mock_browser_frame(slide, l, t, w, h, url="localhost:3000"):
    """Browser chrome with URL bar."""
    _rect(slide, l, t, w, h, CARD_BG, CARD_BORDER)
    chrome = _rect(slide, l, t, w, Inches(0.35), RGBColor(0xE2, 0xE8, 0xF0))
    # traffic lights
    for i, col in enumerate([RGBColor(0xEF, 0x44, 0x44), RGBColor(0xEA, 0xB3, 0x08), RGBColor(0x22, 0xC5, 0x5E)]):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.15 + i * 0.18), t + Inches(0.1), Inches(0.12), Inches(0.12))
        c.fill.solid()
        c.fill.fore_color.rgb = col
        c.line.fill.background()
    _rect(slide, l + Inches(0.75), t + Inches(0.08), w - Inches(0.9), Inches(0.22), WHITE, CARD_BORDER)
    _text(slide, l + Inches(0.85), t + Inches(0.09), w - Inches(1.0), Inches(0.2), url, 8, False, TEXT_MUTED)
    return l, t + Inches(0.35), w, h - Inches(0.35)


def mock_search_bar(slide, l, t, w, query="yaourt", active_view="grid"):
    """Search input + view mode pills."""
    _rect(slide, l, t, w, Inches(0.42), WHITE, CARD_BORDER)
    # search icon circle
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.1), t + Inches(0.08), Inches(0.26), Inches(0.26))
    icon.fill.solid()
    icon.fill.fore_color.rgb = EMERALD
    icon.line.fill.background()
    _text(slide, l + Inches(0.45), t + Inches(0.1), w - Inches(0.55), Inches(0.3), query, 12, False, TEXT_DARK)

    views = [("Grille", "grid"), ("Table", "table"), ("Stats", "stats")]
    x = l + w - Inches(2.4)
    for label, vid in views:
        active = vid == active_view
        col = EMERALD if active else BG_PANEL
        txt = WHITE if active else TEXT_MUTED
        pw = Inches(0.72)
        _rect(slide, x, t + Inches(0.06), pw, Inches(0.3), col, CARD_BORDER if not active else None)
        _text(slide, x, t + Inches(0.1), pw, Inches(0.25), label, 9, active, txt, PP_ALIGN.CENTER)
        x += Inches(0.78)


def mock_product_card(slide, l, t, w, h, name, nutriscore="A", nova="NOVA 1", brand="Danone", variant=0):
    """NFT-style product card mock."""
    grad_top, grad_bot = CARD_GRADIENTS[variant % len(CARD_GRADIENTS)]
    _rect(slide, l, t, w, h, grad_top, EMERALD)
    # image placeholder
    img_h = h * 0.45
    _rect(slide, l + Inches(0.08), t + Inches(0.08), w - Inches(0.16), img_h, grad_bot, CARD_BORDER)
    _text(slide, l + Inches(0.15), t + img_h * 0.3, w - Inches(0.3), Inches(0.3), "image", 8, False, TEXT_MUTED, PP_ALIGN.CENTER)

    y_info = t + img_h + Inches(0.15)
    _text(slide, l + Inches(0.1), y_info, w - Inches(0.2), Inches(0.35), name, 10, True, TEXT_DARK)
    _text(slide, l + Inches(0.1), y_info + Inches(0.32), w - Inches(0.2), Inches(0.25), brand, 8, False, TEXT_MUTED)

    # Nutri-Score badge
    ns_col = NUTRI.get(nutriscore.upper(), EMERALD)
    badge_w = Inches(0.38)
    _rect(slide, l + Inches(0.1), y_info + Inches(0.55), badge_w, Inches(0.28), ns_col)
    _text(slide, l + Inches(0.1), y_info + Inches(0.57), badge_w, Inches(0.24), nutriscore.upper(), 9, True, WHITE, PP_ALIGN.CENTER)

    # NOVA badge
    _rect(slide, l + Inches(0.55), y_info + Inches(0.55), Inches(0.65), Inches(0.28), RGBColor(0x22, 0xC5, 0x5E))
    _text(slide, l + Inches(0.55), y_info + Inches(0.57), Inches(0.65), Inches(0.24), nova, 7, True, WHITE, PP_ALIGN.CENTER)


def mock_filters_panel(slide, l, t, w, h):
    """Sidebar filters mock."""
    _rect(slide, l, t, w, h, BG_PANEL, CARD_BORDER)
    _text(slide, l + Inches(0.12), t + Inches(0.12), w - Inches(0.24), Inches(0.3), "Filtres", 11, True, TEXT_DARK)

    filters = [
        ("Pays", "France"),
        ("Marque", "Danone"),
        ("Nutri-Score", "A, B"),
        ("NOVA", "1, 2"),
    ]
    y = t + Inches(0.45)
    for label, val in filters:
        _text(slide, l + Inches(0.12), y, w - Inches(0.24), Inches(0.2), label, 8, True, TEXT_MUTED)
        _rect(slide, l + Inches(0.12), y + Inches(0.2), w - Inches(0.24), Inches(0.28), WHITE, CARD_BORDER)
        _text(slide, l + Inches(0.18), y + Inches(0.24), w - Inches(0.3), Inches(0.22), val, 9, False, TEXT_DARK)
        y += Inches(0.58)

    # Nutri pills
    _text(slide, l + Inches(0.12), y, w - Inches(0.24), Inches(0.2), "Nutri-Score", 8, True, TEXT_MUTED)
    x = l + Inches(0.12)
    for grade in "ABCDE":
        gw = Inches(0.32)
        _rect(slide, x, y + Inches(0.22), gw, Inches(0.26), NUTRI[grade] if grade in "AB" else RGBColor(0xE2, 0xE8, 0xF0))
        _text(slide, x, y + Inches(0.24), gw, Inches(0.22), grade, 8, True,
              WHITE if grade in "AB" else TEXT_MUTED, PP_ALIGN.CENTER)
        x += Inches(0.36)


def mock_data_table(slide, l, t, w, h):
    """Compact data table with colored Nutri-Score cells."""
    rows = [
        ("Produit", "Nutri", "Énergie", "Sucres"),
        ("Yaourt nature", "A", "61 kcal", "4.5 g"),
        ("Yaourt aux fruits", "B", "98 kcal", "12 g"),
        ("Yaourt 0%", "A", "45 kcal", "5 g"),
    ]
    col_w = [w * 0.4, w * 0.15, w * 0.22, w * 0.23]
    row_h = h / len(rows)

    _rect(slide, l, t, w, h, WHITE, CARD_BORDER)
    y = t
    for ri, row in enumerate(rows):
        x = l
        bg = BG_PANEL if ri == 0 else WHITE
        _rect(slide, l, y, w, row_h, bg)
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            if ri > 0 and ci == 1:
                col = NUTRI.get(cell, EMERALD)
                _rect(slide, x + Inches(0.05), y + Inches(0.06), Inches(0.28), row_h - Inches(0.12), col)
                _text(slide, x + Inches(0.05), y + Inches(0.1), Inches(0.28), row_h - Inches(0.15),
                      cell, 9, True, WHITE, PP_ALIGN.CENTER)
            else:
                bold = ri == 0
                col_txt = TEXT_DARK if ri == 0 else TEXT_MUTED
                _text(slide, x + Inches(0.08), y + Inches(0.08), cw - Inches(0.1), row_h - Inches(0.1),
                      cell, 9 if ri == 0 else 10, bold, col_txt)
            x += cw
        y += row_h


def mock_chart_bars(slide, l, t, w, h):
    """Bar chart Nutri-Score distribution."""
    _rect(slide, l, t, w, h, WHITE, CARD_BORDER)
    _text(slide, l + Inches(0.1), t + Inches(0.08), w - Inches(0.2), Inches(0.25),
          "Distribution Nutri-Score", 10, True, TEXT_DARK)

    data = [("A", 12, "A"), ("B", 8, "B"), ("C", 5, "C"), ("D", 2, "D"), ("E", 1, "E")]
    max_val = max(v for _, v, _ in data)
    bar_area_h = h - Inches(0.55)
    bar_w = (w - Inches(0.4)) / len(data)
    base_y = t + h - Inches(0.25)

    for i, (_, val, grade) in enumerate(data):
        bar_h = bar_area_h * (val / max_val)
        bx = l + Inches(0.2) + i * bar_w
        by = base_y - bar_h
        _rect(slide, bx, by, bar_w - Inches(0.08), bar_h, NUTRI[grade], radius_shape=MSO_SHAPE.RECTANGLE)
        _text(slide, bx, base_y + Inches(0.02), bar_w - Inches(0.08), Inches(0.2),
              grade, 9, True, TEXT_MUTED, PP_ALIGN.CENTER)


def mock_architecture_diagram(slide, l, t, w, h):
    """3-layer architecture with arrows."""
    box_w = (w - Inches(0.6)) / 3
    box_h = h - Inches(0.3)
    layers = [
        ("Frontend", "Next.js", ":3000", EMERALD),
        ("Backend", "Scala/http4s", ":8080", CYAN),
        ("Open Data", "OpenFoodFacts", "externe", VIOLET),
    ]
    x = l
    for i, (title, sub, port, col) in enumerate(layers):
        _rect(slide, x, t, box_w, box_h, WHITE, CARD_BORDER)
        _rect(slide, x, t, box_w, Inches(0.06), col)
        _text(slide, x + Inches(0.1), t + Inches(0.15), box_w - Inches(0.2), Inches(0.3), title, 11, True, TEXT_DARK)
        _text(slide, x + Inches(0.1), t + Inches(0.42), box_w - Inches(0.2), Inches(0.25), sub, 9, False, col)
        _text(slide, x + Inches(0.1), t + box_h - Inches(0.35), box_w - Inches(0.2), Inches(0.25), port, 8, True, TEXT_MUTED, PP_ALIGN.RIGHT)
        if i < 2:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.05), t + box_h / 2 - Inches(0.12),
                Inches(0.25), Inches(0.24)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = TEXT_MUTED
            arr.line.fill.background()
        x += box_w + Inches(0.3)


def mock_json_snippet(slide, l, t, w, h, lines):
    """Dark code block for API response."""
    _rect(slide, l, t, w, h, BG_CODE)
    _rect(slide, l, t, w, Inches(0.04), EMERALD)
    y = t + Inches(0.12)
    for i, line in enumerate(lines):
        col = CODE_ACCENT if i == 0 or "{" in line or "}" in line else CODE_TEXT
        _text(slide, l + Inches(0.15), y, w - Inches(0.3), Inches(0.22), line, 9, False, col, font="Consolas")
        y += Inches(0.22)


def mock_terminal_block(slide, l, t, w, h, lines):
    """Terminal-style command block for demo slide."""
    _rect(slide, l, t, w, h, RGBColor(0x1E, 0x29, 0x3B))
    _rect(slide, l, t, w, Inches(0.04), EMERALD)
    y = t + Inches(0.15)
    for line in lines:
        color = EMERALD if line.startswith("$") else CODE_TEXT
        _text(slide, l + Inches(0.15), y, w - Inches(0.3), Inches(0.22), line, 9, False, color, font="Consolas")
        y += Inches(0.24)
